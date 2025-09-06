from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.conf import settings
from django.http import HttpResponseForbidden, HttpResponse
from django.db.models import Q
from django.views.decorators.csrf import csrf_exempt
from . import models as ds_models 
from sheet_models.models import Analysis, Prompt, RawData
from .forms import *
from smartsheet_sync.fast_apis_side import get_generated_code
from .subprocess import *
from core_performance import performance_logger
from functools import wraps
from bs4 import BeautifulSoup
from django.utils.safestring import mark_safe
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.pdfgen import canvas
import json 
import base64
from io import BytesIO
from weasyprint import HTML, CSS
from concurrent.futures import ThreadPoolExecutor, as_completed
import traceback
import io
import requests
import csv
from sheet_models.data_collector import GetData
from sheet_models.views import handle_uploaded_file
import logging

logger = logging.getLogger(__name__)

@login_required
@performance_logger
def dashboard(request, analysis_id):
    analysis = get_object_or_404(Analysis, id=analysis_id)

    if analysis.prompt.user != request.user:
        return HttpResponseForbidden("Not allowed BadGate")

    data_file = get_object_or_404(RawData, prompt=analysis.prompt)
    data_path = data_file.data.path

    # --- defaults / safe initializations (avoid UnboundLocalError) ---
    desc_data = {}
    corr_data = {}
    quality_data = []
    cols_desc = []
    columns = []
    rows = []
    total_nulls = 0
    total_duplicates = 0

    chart_count = 5
    sample_size = 1000

    plotly_output = []
    plotly_error = ""

    # These are outputs from code execution (separate for prompt vs user)
    prompt_code_output = None
    prompt_code_error = None
    user_code_output = None
    user_code_error = None
    formatted_code_output = []
    user_code_out_raw = ""  # raw marker contents if needed

    check_add_plot = False
    check_more_plots = False

    # --- Load dataframe metadata safely ---
    try:
        ext = handle_uploaded_file(data_file.data)
        df_providor = GetData(data_path, ext)
        raw_df = df_providor.get_loaded_data
        data_shape = raw_df.shape
        max_cols = data_shape[1]
        max_rows = data_shape[0]
    except Exception as e:
        logger.error(f"get_raw_df error {e}")
        raw_df = None
        max_cols = 0
        max_rows = 0

    # --- Safe fetch relations ---
    try:
        desc_obj = ds_models.DescriptiveStatistics.objects.filter(analysis=analysis).first()
        corr_obj = ds_models.Correlation.objects.filter(analysis=analysis).first()
        quality_obj = ds_models.QualityReport.objects.filter(analysis=analysis).first()
        cols_desc_obj = ds_models.ColumnDescription.objects.filter(analysis=analysis).first()
    except Exception as e:
        logger.error(f"Error fetching DS relations: {e}")
        desc_obj = corr_obj = quality_obj = cols_desc_obj = None

    def safe_json_load(data, fallback):
        if not data:
            return fallback
        if isinstance(data, str):
            try:
                return json.loads(data)
            except json.JSONDecodeError:
                return fallback
        return data or fallback

    desc_data = safe_json_load(desc_obj.table if desc_obj else None, {})
    corr_data = safe_json_load(corr_obj.table if corr_obj else None, {})
    quality_data = safe_json_load(quality_obj.table if quality_obj else None, [])
    cols_desc = safe_json_load(cols_desc_obj.table if cols_desc_obj else None, [])

    # --- Prepare correlations ---
    if isinstance(corr_data, dict):
        columns = list(corr_data.keys())
        rows = []
        for row_name, row_vals in corr_data.items():
            row_list = [row_vals.get(col, None) for col in columns]
            rows.append((row_name, row_list))
    else:
        columns = []
        rows = []

    total_nulls = sum(item.get("nulls", 0) for item in quality_data if isinstance(item, dict))
    total_duplicates = sum(item.get("duplicates", 0) for item in quality_data if isinstance(item, dict))

    # --- Forms (guard any form errors) ---
    try:
        form = CodePromptForm(request.POST or None)
    except Exception as e:
        logger.error(f"AI form error: {e}")
        form = CodePromptForm()

    try:
        code_form = UserCodeForm(request.POST or None)
    except Exception as e:
        logger.error(f"UserCode form error: {e}")
        code_form = UserCodeForm()

    try:
        plotly_form = PlotlyChartForm(request.POST or None)
        # Only set choices if raw_df is available and columns detected
        if raw_df is not None and columns:
            try:
                plotly_form.fields["col"].choices = [
                    (col, f"column: {col}, type: {raw_df[col].dtype}") for col in columns
                ]
            except Exception:
                # If there is an issue accessing raw_df types, skip
                pass
    except Exception as e:
        logger.error(f"Plotly form error: {e}")
        plotly_form = PlotlyChartForm()

    try:
        chart_form = ChartSettingsForm(request.POST or None)
        # adjust max values only if we have realistic numbers
        if max_cols:
            try:
                chart_form.fields["chart_count"].max_value = max_cols
            except Exception:
                pass
        if max_rows:
            try:
                chart_form.fields["sample_size"].max_value = max_rows
            except Exception:
                pass
    except Exception as e:
        logger.error(f"ChartSettings form error: {e}")
        chart_form = ChartSettingsForm()

    # --- Handle POST / run tasks ---
    try:
        if request.method == "POST":
            futures = {}
            with ThreadPoolExecutor() as executor:
                # Chart settings
                if chart_form.is_valid():
                    try:
                        chart_count = chart_form.cleaned_data.get("chart_count", chart_count)
                        sample_size = chart_form.cleaned_data.get("sample_size", sample_size)
                        check_more_plots = True
                    except Exception:
                        pass

                # Prompt case (generate code via LLM)
                if form.is_valid():
                    prompt = form.cleaned_data.get("prompt", "")
                    futures["prompt_code"] = executor.submit(
                        process_code_prompt, prompt, data_path, columns
                    )

                # User code case
                if code_form.is_valid():
                    user_code = code_form.cleaned_data.get("user_code", "")
                    futures["user_code"] = executor.submit(
                        process_user_code_secure, user_code, data_path
                    )
                    formatted_code_output.append(str(user_code).strip())

                # Plotly charts (either form specific or default request)
                try:
                    if plotly_form.is_valid():
                        kind = plotly_form.cleaned_data.get("kind")
                        col = plotly_form.cleaned_data.get("col")
                        colors = plotly_form.cleaned_data.get("colors")
                        futures["plotly"] = executor.submit(
                            process_plotly_prompt,
                            "",
                            data_path,
                            columns,
                            {"kind": kind, "col": col, "colors": colors},
                            chart_count,
                            sample_size
                        )
                    else:
                        futures["plotly"] = executor.submit(
                            process_plotly_prompt,
                            "",
                            data_path,
                            columns,
                            {"chart_count": chart_count, "sample_size": sample_size},
                            chart_count,
                            sample_size
                        )
                except Exception as e:
                    logger.error(f"Failed scheduling plotly job: {e}")
                    # schedule a safe default plotly job if available
                    futures["plotly"] = executor.submit(
                        process_plotly_prompt,
                        "",
                        data_path,
                        columns,
                        {"chart_count": chart_count, "sample_size": sample_size},
                        chart_count,
                        sample_size
                    )

                # Wait & collect (with safe handling)
                results = {}
                for k, fut in futures.items():
                    try:
                        results[k] = fut.result()
                    except Exception as e:
                        logger.error(f"Task {k} failed: {e}")
                        results[k] = {"code_output": [], "code_error": str(e), "success": False}

            # --- Collect results safely ---
            # Prompt (LLM) result
            if "prompt_code" in results and results["prompt_code"] is not None:
                try:
                    r = results["prompt_code"]
                    prompt_code_output = r.get("code_output")
                    prompt_code_error = r.get("code_error")
                    formatted_code_output.append(str(r.get("raw_code", "")).strip())
                except Exception as e:
                    logger.error(f"Error parsing prompt_code result: {e}")
                    prompt_code_output = None
                    prompt_code_error = str(e)

            # User code result
            if "user_code" in results and results["user_code"] is not None:
                try:
                    r = results["user_code"]
                    user_code_output = r.get("code_output")
                    user_code_error = r.get("code_error")
                    # also capture raw user section if provided
                    try:
                        user_code_out_raw = r.get("user_code_out") or ""
                    except Exception:
                        user_code_out_raw = ""
                except Exception as e:
                    logger.error(f"Error parsing user_code result: {e}")
                    user_code_output = None
                    user_code_error = str(e)

            # Plotly result
            if "plotly" in results and results["plotly"] is not None:
                try:
                    r = results["plotly"]
                    plotly_output = [mark_safe(json.dumps(fig)) for fig in r.get("plotly_output", [])]
                    plotly_error = r.get("code_error", "") or ""
                    check_add_plot = True
                except Exception as e:
                    logger.error(f"Error parsing plotly result: {e}")
                    plotly_output = []
                    plotly_error = str(e)
                    check_add_plot = False

        else:
            # --- Default charts on first load (GET) ---
            try:
                r = process_plotly_prompt(
                    "", data_path, columns, {"chart_count": chart_count, "sample_size": sample_size}
                )
                plotly_output = [mark_safe(json.dumps(fig)) for fig in r.get("plotly_output", [])]
                plotly_error = r.get("code_error", "") or ""
            except Exception as e:
                logger.error(f"Default chart error: {e}")
                plotly_output, plotly_error = [], ""
    except Exception as e:
        logger.error(f"Error while handling forms: {e}\n{traceback.format_exc()}")

    # --- Prepare combined fallback outputs for templates ---
    # primary outputs (prefer prompt then user for "main" display if both exist)
    main_code_output = prompt_code_output if prompt_code_output else user_code_output
    main_code_error = prompt_code_error if prompt_code_error else user_code_error

    # For raw user_code_out field used by template (previously user_code_out: code_output or "")
    user_code_out_field = user_code_out_raw or (user_code_output if isinstance(user_code_output, str) else "")

    return render(
        request,
        settings.T_PATH["dashboard"],
        {   
            "analysis_id": analysis_id,
            "corr_columns": columns,
            "corr_rows": rows,
            "desc": desc_data,
            "quality": quality_data,
            "total_nulls": total_nulls,
            "total_duplicates": total_duplicates,
            "title": f"DashBoard {analysis_id}",
            "data_title": analysis.title,
            "data_description": analysis.description,
            "cols_desc": cols_desc,
            "form": form,
            "code_form": code_form,
            "plotly_form": plotly_form,
            "chart_form": chart_form,
            "chart_count": chart_count,
            "sample_size": sample_size,
            # separated outputs
            "prompt_code_output": prompt_code_output,
            "prompt_code_error": prompt_code_error,
            "user_code_output": user_code_output,
            "user_code_error": user_code_error,
            # main combined fallback used by older template pieces
            "code_output": main_code_output,
            "code_error": main_code_error,
            "generated_code": formatted_code_output,
            "plotly_output": plotly_output,
            "plotly_error": plotly_error,
            "user_code_out": user_code_out_field or "",
            "check_add_plot": check_add_plot,
            "check_more_plots": check_more_plots,
        },
    )
    
@login_required
@performance_logger
def dashboard_history(request): 
    prompts = Prompt.objects.filter(user=request.user)
    analyses = Analysis.objects.filter(
        prompt__in=prompts
    ).filter(
        Q(descriptive_statistics__isnull=False) |
        Q(correlations__isnull=False) |
        Q(quality_reports__isnull=False) |
        Q(column_desciptions__isnull=False)
    ).distinct().order_by('-created_at')

    return render(request, settings.T_PATH["history"], {
        "analyses": analyses,
        "title": "History",
    })
    
@login_required
@performance_logger
def generate_code_view(request):
    code_result = None
    if request.method == "POST":
        form = CodePromptForm(request.POST)
        if form.is_valid():
            prompt = form.cleaned_data["prompt"]
            code_result = get_generated_code(prompt).get("code")
            print(code_result)
    else:
        form = CodePromptForm()

    return render(request, settings.T_PATH["code_prompt"], {"form": form, "code_result": code_result})

# @csrf_exempt
# @performance_logger
# def export_pdf(request):
#     if request.method == "POST":
#         data = json.loads(request.body.decode("utf-8"))
#         html_content = data.get("html", "")

#         # Add custom CSS for print styling
#         css = CSS(string="""
#             body { font-family: 'Arial', sans-serif; color: #333; }
#             h1, h2, h3 { color: #2c3e50; margin-top: 20px; }
#             table { width: 100%; border-collapse: collapse; margin-bottom: 20px; }
#             th, td { border: 1px solid #ddd; padding: 8px; }
#             th { background: #f4f6f8; text-align: left; }
#             img { max-width: 100%; border-radius: 6px; margin: 10px 0; }
#             p { font-size: 14px; line-height: 1.5; }
#         """)

#         pdf_file = HTML(string=html_content).write_pdf(stylesheets=[css])

#         response = HttpResponse(pdf_file, content_type="application/pdf")
#         response["Content-Disposition"] = 'attachment; filename="report.pdf"'
#         return response

#     return HttpResponse("Method not allowed", status=405)

# @csrf_exempt
# @performance_logger
# def export_pptx(request):
#     if request.method == "POST":
#         data = json.loads(request.body)
#         html_content = data.get("html", "")
#         soup = BeautifulSoup(html_content, "html.parser")

#         prs = Presentation()
#         blank_layout = prs.slide_layouts[5]

#         for elem in soup.find_all(["h1", "h2", "h3", "p", "table", "img"]):
#             slide = prs.slides.add_slide(blank_layout)

#             if elem.name in ["h1", "h2", "h3", "p"]:
#                 txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
#                 tf = txBox.text_frame
#                 p = tf.add_paragraph()
#                 p.text = elem.get_text().strip()

#                 # Styling
#                 if elem.name == "h1":
#                     p.font.size = Pt(32)
#                     p.font.bold = True
#                     p.font.color.rgb = RGBColor(44, 62, 80)
#                 elif elem.name == "h2":
#                     p.font.size = Pt(26)
#                     p.font.color.rgb = RGBColor(52, 73, 94)
#                 elif elem.name == "h3":
#                     p.font.size = Pt(22)
#                     p.font.color.rgb = RGBColor(127, 140, 141)
#                 else:  # p
#                     p.font.size = Pt(16)
#                     p.font.color.rgb = RGBColor(80, 80, 80)

#             elif elem.name == "table":
#                 rows = elem.find_all("tr")
#                 cols = rows[0].find_all(["td", "th"])
#                 table = slide.shapes.add_table(len(rows), len(cols), Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table

#                 # Table styling
#                 for r, row in enumerate(rows):
#                     for c, cell in enumerate(row.find_all(["td", "th"])):
#                         table.cell(r, c).text = cell.get_text().strip()
#                         cell_text = table.cell(r, c).text_frame.paragraphs[0]
#                         cell_text.font.size = Pt(14)
#                         if r == 0:  # header row
#                             cell_text.font.bold = True
#                             cell_text.font.color.rgb = RGBColor(255, 255, 255)
#                             table.cell(r, c).fill.solid()
#                             table.cell(r, c).fill.fore_color.rgb = RGBColor(52, 73, 94)

#             elif elem.name == "img":
#                 img_url = elem["src"]
#                 try:
#                     if img_url.startswith("data:image"):  # base64
#                         header, encoded = img_url.split(",", 1)
#                         img_data = base64.b64decode(encoded)
#                     else:
#                         img_data = requests.get(img_url).content

#                     img_stream = io.BytesIO(img_data)
#                     slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), height=Inches(4))
#                 except Exception:
#                     pass

#         buffer = io.BytesIO()
#         prs.save(buffer)
#         buffer.seek(0)

#         response = HttpResponse(
#             buffer,
#             content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
#         )
#         response["Content-Disposition"] = 'attachment; filename="report.pptx"'
#         return response

#     return HttpResponse(status=405)

@csrf_exempt
@performance_logger
def download_csv(request):
    """
    Convert the HTML table sent from frontend into CSV.
    Expects POST with 'html' field containing <table>...</table>.
    """
    if request.method == "POST":
        html = request.POST.get("html", "")
        if not html:
            return HttpResponse("No table provided", status=400)

        soup = BeautifulSoup(html, "html.parser")
        table = soup.find("table")
        if not table:
            return HttpResponse("No table found in HTML", status=400)

        # Create HTTP response with CSV content
        response = HttpResponse(content_type="text/csv")
        response["Content-Disposition"] = 'attachment; filename="table.csv"'

        writer = csv.writer(response)

        for row in table.find_all("tr"):
            cols = [c.get_text(strip=True) for c in row.find_all(["td", "th"])]
            writer.writerow(cols)

        return response

    return HttpResponse("Only POST allowed", status=405)


@csrf_exempt
@performance_logger
def export_pdf(request):
    if request.method == "POST":
        data = json.loads(request.body.decode("utf-8"))
        html_content = data.get("html", "")

        # Add custom CSS to ensure all scrollable content is visible
        css = CSS(string="""
            body { 
                font-family: 'Arial', sans-serif; 
                color: #333; 
                margin: 20px;
            }
            h1, h2, h3 { 
                color: #2c3e50; 
                margin-top: 20px; 
            }
            table { 
                width: 100%; 
                border-collapse: collapse; 
                margin-bottom: 20px; 
                page-break-inside: auto;
            }
            th, td { 
                border: 1px solid #ddd; 
                padding: 8px; 
                text-align: left;
            }
            th { 
                background: #f4f6f8; 
                font-weight: bold; 
            }
            img { 
                max-width: 100%; 
                border-radius: 6px; 
                margin: 10px 0; 
            }
            p { 
                font-size: 14px; 
                line-height: 1.5; 
            }
            ul, ol { 
                font-size: 14px; 
                line-height: 1.5; 
                margin: 10px 0; 
                padding-left: 20px;
            }
            /* Ensure scrollable elements are fully visible */
            div, table, ul, ol { 
                overflow: visible !important; 
                max-height: none !important; 
                height: auto !important; 
                page-break-inside: auto;
            }
            /* Handle large tables */
            table tr { 
                page-break-inside: avoid; 
                page-break-after: auto;
            }
        """)

        # Parse HTML to ensure all content is included
        soup = BeautifulSoup(html_content, "html.parser")
        
        # Remove inline styles that might cause scrolling
        for elem in soup.find_all(['div', 'table', 'ul', 'ol']):
            if 'style' in elem.attrs:
                style = elem.attrs['style']
                # Remove overflow and max-height properties
                new_style = ';'.join([s for s in style.split(';') 
                                    if not any(x in s.lower() for x in ['overflow', 'max-height'])])
                if new_style.strip():
                    elem.attrs['style'] = new_style
                else:
                    del elem.attrs['style']

        # Convert modified HTML back to string
        modified_html = str(soup)

        # Generate PDF with modified HTML
        pdf_file = HTML(string=modified_html).write_pdf(stylesheets=[css])

        response = HttpResponse(pdf_file, content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename="report.pdf"'
        return response

    return HttpResponse("Method not allowed", status=405)

@csrf_exempt
@performance_logger
def export_pptx(request):
    if request.method == "POST":
        data = json.loads(request.body)
        html_content = data.get("html", "")
        soup = BeautifulSoup(html_content, "html.parser")

        prs = Presentation()
        blank_layout = prs.slide_layouts[5]

        for elem in soup.find_all(["h1", "h2", "h3", "p", "table", "ul", "ol", "img"]):
            slide = prs.slides.add_slide(blank_layout)

            if elem.name in ["h1", "h2", "h3", "p"]:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = elem.get_text().strip()

                # Styling
                if elem.name == "h1":
                    p.font.size = Pt(32)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(44, 62, 80)
                elif elem.name == "h2":
                    p.font.size = Pt(26)
                    p.font.color.rgb = RGBColor(52, 73, 94)
                elif elem.name == "h3":
                    p.font.size = Pt(22)
                    p.font.color.rgb = RGBColor(127, 140, 141)
                else:  # p
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(80, 80, 80)

            elif elem.name == "table":
                rows = elem.find_all("tr")
                cols = rows[0].find_all(["td", "th"])
                table = slide.shapes.add_table(len(rows), len(cols), Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table

                # Table styling
                for r, row in enumerate(rows):
                    for c, cell in enumerate(row.find_all(["td", "th"])):
                        table.cell(r, c).text = cell.get_text().strip()
                        cell_text = table.cell(r, c).text_frame.paragraphs[0]
                        cell_text.font.size = Pt(14)
                        if r == 0:  # header row
                            cell_text.font.bold = True
                            cell_text.font.color.rgb = RGBColor(255, 255, 255)
                            table.cell(r, c).fill.solid()
                            table.cell(r, c).fill.fore_color.rgb = RGBColor(52, 73, 94)

            elif elem.name in ["ul", "ol"]:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                tf = txBox.text_frame
                for li in elem.find_all("li"):
                    p = tf.add_paragraph()
                    p.text = li.get_text().strip()
                    p.level = 0 if elem.name == "ul" else 1  # Bullet or numbered list
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(80, 80, 80)

            elif elem.name == "img":
                img_url = elem["src"]
                try:
                    if img_url.startswith("data:image"):  # base64
                        header, encoded = img_url.split(",", 1)
                        img_data = base64.b64decode(encoded)
                    else:
                        img_data = requests.get(img_url).content

                    img_stream = io.BytesIO(img_data)
                    slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), height=Inches(4))
                except Exception:
                    pass

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        response = HttpResponse(
            buffer,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        response["Content-Disposition"] = 'attachment; filename="report.pptx"'
        return response

    return HttpResponse(status=405)